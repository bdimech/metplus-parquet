"""Build METplus to Parquet summary presentation using the bureau template."""

import io
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd
import pyarrow.parquet as pq

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# -- Colours ---------------------------------------------------------------
BLUE    = RGBColor(0x24, 0x61, 0xE5)
DBLUE   = RGBColor(0x15, 0x3A, 0xBF)
LBLUE   = RGBColor(0xE8, 0xF0, 0xFF)
DARK    = RGBColor(0x1C, 0x1C, 0x28)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGREY   = RGBColor(0xF5, 0xF5, 0xF7)
MGREY   = RGBColor(0xB8, 0xB8, 0xC4)
DGREY   = RGBColor(0x6E, 0x6E, 0x78)
GREEN   = RGBColor(0x10, 0x7C, 0x10)
LGREEN  = RGBColor(0xE6, 0xF5, 0xE6)
ORANGE  = RGBColor(0xC0, 0x48, 0x00)
LORANGE = RGBColor(0xFD, 0xED, 0xE5)
TEAL    = RGBColor(0x00, 0x78, 0x8C)
LTEAL   = RGBColor(0xDE, 0xF4, 0xF7)

def i(v): return Inches(v)
def p(v): return Pt(v)


# -- Template setup --------------------------------------------------------
prs = Presentation("METplus - Application Support.pptx")
sldIdLst = prs.slides._sldIdLst
while len(sldIdLst):
    prs.part.drop_rel(sldIdLst[0].get(qn("r:id")))
    del sldIdLst[0]

LAYOUT = prs.slide_layouts[6]


# -- Helpers ---------------------------------------------------------------
def new_slide():
    """Add a blank slide, stripping all template placeholder shapes."""
    slide = prs.slides.add_slide(LAYOUT)
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    return slide


def _no_line(shp):
    """Set shape border to no-fill (invisible)."""
    spPr = shp._element.spPr
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    for child in list(ln):
        ln.remove(child)
    etree.SubElement(ln, qn("a:noFill"))


def add_title(slide, text):
    tb = slide.shapes.add_textbox(i(0.45), i(0.20), i(12.4), i(0.52))
    tf = tb.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.bold = True
    run.font.size = p(22)
    run.font.color.rgb = BLUE
    rule = slide.shapes.add_shape(1, i(0.45), i(0.76), i(12.4), i(0.022))
    rule.fill.solid()
    rule.fill.fore_color.rgb = BLUE
    _no_line(rule)


def rect(slide, left, top, w, h, text="",
         fill=BLUE, fg=WHITE, border=None, bw=0,
         fs=10, bold=False, align=PP_ALIGN.CENTER, wrap=True, rounded=True):
    shp = slide.shapes.add_shape(5 if rounded else 1, i(left), i(top), i(w), i(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if bw == 0 or border is None:
        _no_line(shp)
    else:
        shp.line.color.rgb = border
        shp.line.width = p(bw)
    tf = shp.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left   = i(0.08)
    tf.margin_right  = i(0.08)
    tf.margin_top    = i(0.04)
    tf.margin_bottom = i(0.04)
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.color.rgb = fg
    run.font.size = p(fs)
    run.font.bold = bold
    return shp


def tbx(slide, left, top, w, h, lines,
        fs=10, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(i(left), i(top), i(w), i(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.size = p(fs)
        run.font.color.rgb = color
        run.font.bold = bold
    return tb


def arrow(slide, x1, y1, x2, y2, color=MGREY, w=1.5):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, i(x1), i(y1), i(x2), i(y2))
    c.line.color.rgb = color
    c.line.width = p(w)
    spPr = c._element.spPr
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        te = etree.SubElement(ln, qn("a:tailEnd"))
        te.set("type", "arrow")
        te.set("w", "med")
        te.set("len", "med")
    return c


# ==========================================================================
# SLIDE 1 -- What We Get: METplus Output Files
# ==========================================================================
s1 = new_slide()
add_title(s1, "What We Get: METplus Verification Outputs")

tbx(s1, 0.45, 0.95, 5.6, 5.5, [
    "METplus runs routine NWP verification and writes",
    "results to structured output folders.",
    "",
    "Each configuration produces a folder named:",
    "  <Model>_<Obs>_<Parameter>_<Timestep>",
    "",
    "Inside each folder are daily YYYYMMDDHH",
    "subdirectories, each containing .stat files",
    "for every forecast lead time and domain.",
    "",
    "Two types of verification output:",
    "",
    "  GridStat - deterministic scores",
    "  (bias, RMSE, MAE)",
    "",
    "  EnsembleStat - ensemble scores",
    "  (CRPS, Spread, Skill Score)",
], fs=11)

# -- Folder tree -----------------------------------------------------------
RX = 6.7
# Root box
rect(s1, RX, 1.0, 6.2, 0.42, "input_all/",
     fill=BLUE, fg=WHITE, fs=10, bold=True)

# Connector trunk
arrow(s1, RX + 1.55, 1.42, RX + 1.55, 1.68, color=MGREY, w=1.0)
arrow(s1, RX + 4.65, 1.42, RX + 4.65, 1.68, color=MGREY, w=1.0)
arrow(s1, RX + 1.55, 1.55, RX + 4.65, 1.55, color=MGREY, w=1.0)

# Branch folders
rect(s1, RX,       1.68, 3.0, 0.40,
     "AGlobal4_Analysis_T_AllLevels_00Z/",
     fill=DBLUE, fg=WHITE, fs=8.5, bold=False)
rect(s1, RX + 3.1, 1.68, 3.0, 0.40,
     "AGlobal4E_Analysis_T850_00Z/",
     fill=DBLUE, fg=WHITE, fs=8.5, bold=False)

# Subfolder connectors (left branch)
arrow(s1, RX + 1.5, 2.08, RX + 1.5, 2.30, color=MGREY, w=0.8)
# GridStat subfolder (left)
rect(s1, RX, 2.30, 3.0, 0.36, "GridStat/  2026041000/",
     fill=LGREY, fg=DGREY, fs=8.5, border=MGREY, bw=0.4, rounded=False)

# File type pills (left branch)
rect(s1, RX,      2.80, 0.92, 0.28, "_cnt.txt",
     fill=GREEN, fg=WHITE, fs=8)
rect(s1, RX+0.97, 2.80, 0.98, 0.28, "_sl1l2.txt",
     fill=GREEN, fg=WHITE, fs=8)
rect(s1, RX+2.0,  2.80, 0.98, 0.28, "_sal1l2.txt",
     fill=GREEN, fg=WHITE, fs=8)

# Subfolder connectors (right branch)
arrow(s1, RX + 4.6, 2.08, RX + 4.6, 2.30, color=MGREY, w=0.8)
# GridStat + EnsembleStat subfolders (right branch)
rect(s1, RX+3.1,  2.30, 1.45, 0.36, "GridStat/",
     fill=LGREY, fg=DGREY, fs=8.5, border=MGREY, bw=0.4, rounded=False)
rect(s1, RX+4.65, 2.30, 1.45, 0.36, "EnsembleStat/",
     fill=LGREY, fg=DGREY, fs=8.5, border=MGREY, bw=0.4, rounded=False)
rect(s1, RX+3.1,  2.80, 1.45, 0.28, "_cnt.txt",
     fill=GREEN, fg=WHITE, fs=8)
rect(s1, RX+4.65, 2.80, 1.45, 0.28, "_ecnt.txt",
     fill=ORANGE, fg=WHITE, fs=8)

# Legend
tbx(s1, RX, 3.30, 6.2, 0.28, ["Legend"], fs=8.5, color=DGREY, bold=True)
rect(s1, RX,     3.60, 0.95, 0.28, "GridStat",
     fill=GREEN, fg=WHITE, fs=8)
tbx(s1, RX+1.05, 3.60, 2.2, 0.30,
    ["Deterministic (CNT, SL1L2, SAL1L2)"], fs=8.5, color=DARK)
rect(s1, RX+3.1, 3.60, 1.2, 0.28, "EnsembleStat",
     fill=ORANGE, fg=WHITE, fs=8)
tbx(s1, RX+4.4, 3.60, 2.0, 0.30,
    ["Ensemble (ECNT)"], fs=8.5, color=DARK)


# ==========================================================================
# SLIDE 2 -- Why Parquet?
# ==========================================================================
s2 = new_slide()
add_title(s2, "Why Parquet?")

tbx(s2, 0.45, 0.93, 12.4, 0.38,
    ["Moving from hundreds of daily .stat files to a single, portable, queryable dataset per model."],
    fs=11, color=DGREY)

# 4 benefit cards
CARDS = [
    ("Columnar\nStorage", BLUE,
     "Reads only the columns you request."
     " Skip irrelevant fields entirely - no scanning of full rows."),
    ("Compressed\nby Default", TEAL,
     "5-10x smaller than CSV. Snappy compression is applied automatically."
     " Fits on a USB drive."),
    ("Self-\nDescribing", GREEN,
     "Schema, data types, and metadata all embedded in one file."
     " No separate DDL or config needed."),
    ("Python\nNative", DBLUE,
     "One line to load with pandas."
     " Filter columns before reading with pyarrow - no SQL needed."),
]

CW, CH  = 2.8, 1.05
CSEP    = 0.22
CX0     = 0.55
CT      = 1.48

for ci, (title, color, desc) in enumerate(CARDS):
    cx = CX0 + ci * (CW + CSEP)
    rect(s2, cx, CT, CW, CH, title,
         fill=color, fg=WHITE, fs=12, bold=True)
    tbx(s2, cx, CT + CH + 0.12, CW, 1.55,
        [desc], fs=10, color=DARK)

# vs SQL comparison strip
VS_TEXT = (
    "vs a SQL database:  no server to install, no connection strings, no schema migrations, "
    "no DBA required.  Share your verification data as a single .parquet file - "
    "it opens in one line of Python on any machine."
)
rect(s2, 0.45, 4.20, 12.4, 0.72, VS_TEXT,
     fill=LGREY, fg=DGREY, fs=10, rounded=False, align=PP_ALIGN.LEFT)

# One-liner code teaser
tbx(s2, 0.55, 5.06, 12.0, 0.35,
    ['  df = pq.read_table("AGlobal4_Analysis.parquet", '
     'filters=[(\"PARAMETER\",\"=\",\"MSLP\")]).to_pandas()'],
    fs=10, color=BLUE)


# ==========================================================================
# SLIDE 3 -- Conversion Pipeline  (was slide 2)
# ==========================================================================
s3 = new_slide()
add_title(s3, "What We Do: Conversion Pipeline")

BW, BH = 2.18, 0.72
FY = 2.2
XS = [0.55 + k * (BW + 0.20) for k in range(5)]

PIPE_BOXES = [
    ("Input\nFolders",      BLUE,  WHITE),
    ("Read &\nParse",       BLUE,  WHITE),
    ("Merge &\nClean",      BLUE,  WHITE),
    ("Group by\nModel/Obs", BLUE,  WHITE),
    ("Parquet\nFiles",      GREEN, WHITE),
]
for k, (label, fill, fg) in enumerate(PIPE_BOXES):
    rect(s3, XS[k], FY, BW, BH, label,
         fill=fill, fg=fg, fs=11, bold=True)
    if k < 4:
        arrow(s3, XS[k] + BW, FY + BH/2,
              XS[k+1],         FY + BH/2,
              color=MGREY, w=1.5)

PIPE_SUBS = [
    ["<Model>_<Obs>_<Param>_<TS>",
     "GridStat/ or EnsembleStat/",
     "YYYYMMDDHH subdirectories"],
    ["CNT, SL1L2, SAL1L2",
     "(GridStat)",
     "ECNT  (EnsembleStat)"],
    ["Drop metadata cols",
     "Convert lead to hours",
     "Add PARAMETER column",
     "Numeric type casting"],
    ["One output per",
     "model + obs pair:",
     "AGlobal4_Analysis",
     "AGlobal4E_Analysis"],
    ["One .parquet per",
     "model/obs pair",
     "All parameters",
     "All timesteps"],
]
for k, sub in enumerate(PIPE_SUBS):
    tbx(s3, XS[k], FY + BH + 0.14, BW, 1.6,
        sub, fs=9.5, color=DGREY, align=PP_ALIGN.CENTER)

# Safe re-run note
rect(s3, 1.8, 5.08, 9.7, 0.52,
     "Re-running is safe - already-loaded init dates are skipped per PARAMETER and STAT_TYPE",
     fill=LGREY, fg=DGREY, fs=9.5, rounded=False)


# ==========================================================================
# SLIDE 4 -- Parquet Schema  (was slide 3)
# ==========================================================================
s4 = new_slide()
add_title(s4, "What the Outputs Look Like: Parquet Schema")

tbx(s4, 0.45, 0.92, 12.4, 0.38,
    ["Each row = one unique combination of init date, lead time, pressure level, domain, parameter, and stat type."],
    fs=10.5)

CDEFS = [
    ("PARAMETER",       0.45, 1.65, BLUE,   WHITE, BLUE),
    ("STAT_TYPE",       2.15, 1.55, BLUE,   WHITE, BLUE),
    ("INIT_DATE",       3.75, 1.55, BLUE,   WHITE, BLUE),
    ("FCST_LEAD_H",     5.35, 1.40, BLUE,   WHITE, BLUE),
    ("FCST_LEV",        6.80, 1.35, BLUE,   WHITE, BLUE),
    ("VX_MASK",         8.20, 1.35, BLUE,   WHITE, BLUE),
    ("ME / RMSE / MAE", 9.60, 1.75, GREEN,  WHITE, GREEN),
    ("CRPS / SPREAD",  11.40, 1.80, ORANGE, WHITE, ORANGE),
]
HY = 1.55
RH = 0.32

for lbl, cx, cw, f, fg, b in CDEFS:
    rect(s4, cx, HY, cw, RH, lbl,
         fill=f, fg=fg, fs=8.5, bold=True, rounded=False)

DATA = [
    ("T_AllLevels", "GridStat",     "2026-04-10", "24",  "P850", "Australia", "-0.5 / 1.87 / 1.43", "  -"),
    ("MSLP",        "GridStat",     "2026-04-10", "48",  "L0",   "NH",        "0.3 / 2.1 / 1.8",    "  -"),
    ("T850",        "EnsembleStat", "2026-04-15", "24",  "P850", "NH",        "0.28 / 3.26 / 2.26", "1.62 / 2.73"),
    ("Wind_All",    "GridStat",     "2026-04-15", "72",  "P250", "SH",        "0.1 / 2.4 / 1.9",    "  -"),
]
for ri, row in enumerate(DATA):
    ry = HY + RH + ri * RH
    rbg = WHITE if ri % 2 == 0 else LGREY
    for ci, (lbl, cx, cw, _f, _fg, _b) in enumerate(CDEFS):
        cell_fill = rbg
        cell_fg   = DARK
        if ci == 1 and row[ci] == "EnsembleStat":
            cell_fill = LORANGE
            cell_fg   = ORANGE
        elif ci == 1:
            cell_fill = LGREEN
            cell_fg   = GREEN
        rect(s4, cx, ry, cw, RH, row[ci],
             fill=cell_fill, fg=cell_fg,
             border=RGBColor(0xDD, 0xDD, 0xDD), bw=0.3, fs=8.5, rounded=False)

legend_y = HY + RH * (1 + len(DATA)) + 0.18
rect(s4, 0.45, legend_y, 1.4, 0.28, "Index columns",
     fill=BLUE,   fg=WHITE, fs=8, rounded=False)
rect(s4, 1.95, legend_y, 1.4, 0.28, "GridStat stats",
     fill=GREEN,  fg=WHITE, fs=8, rounded=False)
rect(s4, 3.45, legend_y, 1.6, 0.28, "EnsembleStat stats",
     fill=ORANGE, fg=WHITE, fs=8, rounded=False)

tbx(s4, 0.45, legend_y + 0.38, 12.4, 0.4,
    ["Constant columns (model, variable, units, obs type) are stripped from rows and embedded as file-level metadata."],
    fs=9, color=DGREY)


# ==========================================================================
# SLIDE 5 -- How We Can Use the Data  (was slide 4)
# ==========================================================================
s5 = new_slide()
add_title(s5, "How We Can Use the Data")

USE = [
    ("Single Model\nExploration",
     "Load one model's file, filter to a parameter,\n"
     "level and domain, and plot any metric by lead time.",
     BLUE),
    ("Cross-Model\nComparison",
     "Load two files side-by-side, filter to the same\n"
     "parameter and domain, and compare RMSE or bias.",
     BLUE),
    ("Ensemble\nDiagnostics",
     "Filter STAT_TYPE = EnsembleStat to access CRPS\n"
     "and Spread. Check Spread ~ RMSE for calibration.",
     ORANGE),
]
for ui, (ut, ud, color) in enumerate(USE):
    uy = 1.05 + ui * 1.88
    rect(s5, 0.45, uy, 2.5, 0.62, ut,
         fill=color, fg=WHITE, fs=11, bold=True)
    tbx(s5, 0.45, uy + 0.70, 2.55, 1.1,
        ud.split("\n"), fs=10, color=DARK)

# Code box
CODE_LINES = [
    ("import pyarrow.parquet as pq",               "kw"),
    ("",                                            ""),
    ("# Load a model dataset",                      "cm"),
    ("g4 = pq.read_table(",                         "tx"),
    ("    'AGlobal4_Analysis.parquet'",             "st"),
    (").to_pandas()",                               "tx"),
    ("",                                            ""),
    ("# Filter to one parameter",                   "cm"),
    ("t = g4[g4['PARAMETER'] == 'T_AllLevels']",   "tx"),
    ("",                                            ""),
    ("# Cross-model comparison",                    "cm"),
    ("g4e = pq.read_table(",                        "tx"),
    ("    'AGlobal4E_Analysis.parquet',",           "st"),
    ("    filters=[('PARAMETER','=','MSLP')]",      "tx"),
    (").to_pandas()",                               "tx"),
    ("",                                            ""),
    ("# Ensemble stats only",                       "cm"),
    ("ens = g4e[",                                  "tx"),
    ("    g4e['STAT_TYPE'] == 'EnsembleStat']",     "tx"),
]
COLORS = {
    "kw": RGBColor(0x56, 0x9C, 0xD6),
    "cm": RGBColor(0x6A, 0x99, 0x55),
    "st": RGBColor(0xCE, 0x91, 0x78),
    "tx": RGBColor(0xCE, 0xD1, 0xD8),
    "":   RGBColor(0xCE, 0xD1, 0xD8),
}
code_shp = s5.shapes.add_shape(1, i(3.2), i(1.0), i(9.9), i(5.6))
code_shp.fill.solid()
code_shp.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)
_no_line(code_shp)
tf = code_shp.text_frame
tf.word_wrap = False
tf.margin_left  = i(0.20)
tf.margin_top   = i(0.14)
tf.margin_bottom = i(0.08)
for li, (line, kind) in enumerate(CODE_LINES):
    para = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
    run  = para.add_run()
    run.text = line
    run.font.size = p(9.5)
    run.font.color.rgb = COLORS[kind]


# ==========================================================================
# SLIDE 6 -- Example Results  (was slide 5)
# ==========================================================================
s6 = new_slide()
add_title(s6, "Example Results")

G4_PATH  = Path("output_all/AGlobal4_Analysis.parquet")
G4E_PATH = Path("output_all/AGlobal4E_Analysis.parquet")

if G4_PATH.exists() and G4E_PATH.exists():
    g4  = pq.read_table(G4_PATH).to_pandas()
    g4e = pq.read_table(G4E_PATH).to_pandas()

    # Chart 1: MSLP RMSE by lead -- 12Z inits only (AGlobal4E 00Z is anomaly form)
    fig1, ax1 = plt.subplots(figsize=(5.8, 3.4), dpi=150)
    for label_m, df_m, col_m in [("AGlobal4", g4, "#2461E5"),
                                   ("AGlobal4E", g4e, "#C04800")]:
        sub = df_m[
            (df_m["PARAMETER"] == "MSLP") &
            (df_m["STAT_TYPE"] == "GridStat") &
            (df_m["FCST_LEV"]  == "L0") &
            (df_m["VX_MASK"]   == "NH") &
            (pd.to_datetime(df_m["INIT_DATE"]).dt.hour == 12)
        ]
        if not sub.empty:
            s = sub.groupby("FCST_LEAD_H")["RMSE"].mean()
            ax1.plot(s.index, s.values, marker="o", label=label_m,
                     color=col_m, linewidth=2, markersize=4)
    ax1.set_xlabel("Forecast Lead Time (hours)", fontsize=9)
    ax1.set_ylabel("RMSE (Pa)", fontsize=9)
    ax1.set_title("MSLP - RMSE by Lead Time over NH (12Z inits)",
                  fontsize=10, fontweight="bold")
    ax1.legend(fontsize=8)
    ax1.grid(True, alpha=0.2, linestyle="--")
    ax1.spines[["top", "right"]].set_visible(False)
    fig1.tight_layout()
    buf1 = io.BytesIO()
    fig1.savefig(buf1, format="png", bbox_inches="tight")
    buf1.seek(0)
    plt.close(fig1)

    # Chart 2: T850 Spread vs RMSE -- AGlobal4E EnsembleStat
    fig2, ax2 = plt.subplots(figsize=(5.8, 3.4), dpi=150)
    ens = g4e[
        (g4e["PARAMETER"] == "T850") &
        (g4e["STAT_TYPE"] == "EnsembleStat") &
        (g4e["FCST_LEV"]  == "P850") &
        (g4e["VX_MASK"]   == "NH")
    ]
    if not ens.empty:
        spread = ens.groupby("FCST_LEAD_H")["SPREAD"].mean()
        rmse   = ens.groupby("FCST_LEAD_H")["RMSE"].mean()
        ax2.plot(spread.index, spread.values, marker="o", label="Spread",
                 color="#C04800", linewidth=2, markersize=4)
        ax2.plot(rmse.index, rmse.values, marker="s", label="RMSE",
                 color="#2461E5", linewidth=2, markersize=4)
    ax2.set_xlabel("Forecast Lead Time (hours)", fontsize=9)
    ax2.set_ylabel("Value (K)", fontsize=9)
    ax2.set_title("AGlobal4E T850 - Spread vs RMSE over NH",
                  fontsize=10, fontweight="bold")
    ax2.legend(fontsize=8)
    ax2.grid(True, alpha=0.2, linestyle="--")
    ax2.spines[["top", "right"]].set_visible(False)
    fig2.tight_layout()
    buf2 = io.BytesIO()
    fig2.savefig(buf2, format="png", bbox_inches="tight")
    buf2.seek(0)
    plt.close(fig2)

    s6.shapes.add_picture(buf1, i(0.4),  i(1.05), i(6.1), i(3.85))
    s6.shapes.add_picture(buf2, i(6.85), i(1.05), i(6.1), i(3.85))

    tbx(s6, 0.4,  5.0, 6.1, 0.4,
        ["Cross-model comparison: MSLP RMSE by lead time over NH (12Z inits only)."],
        fs=9, color=DGREY, align=PP_ALIGN.CENTER)
    tbx(s6, 6.85, 5.0, 6.1, 0.4,
        ["Ensemble calibration: Spread ~ RMSE indicates a well-calibrated ensemble."],
        fs=9, color=DGREY, align=PP_ALIGN.CENTER)

    tbx(s6, 0.45, 5.58, 12.4, 0.7, [
        'Load:    pq.read_table("AGlobal4_Analysis.parquet", filters=[("PARAMETER","=","MSLP")]).to_pandas()',
        'Filter:  df[df["STAT_TYPE"] == "EnsembleStat"]  |  df[df["VX_MASK"] == "NH"]',
    ], fs=8.5, color=DGREY)
else:
    tbx(s6, 0.5, 2.5, 12.0, 1.0,
        ["Run convert_all() first to generate output_all/, then re-run make_pptx.py."],
        fs=12, color=DARK)


# -- Save ------------------------------------------------------------------
out = "METplus - Parquet Summary.pptx"
prs.save(out)
print(f"Saved: {out}  ({prs.slides._sldIdLst.__len__()} slides)")
